"""Deterministic renderers for SW-012 public and presenter outputs."""

from __future__ import annotations

import hashlib
import html
import os
import tempfile
from datetime import datetime
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen.canvas import Canvas

from ..visual.grammar import GRAMMAR_SCHEMA_VERSION, html_theme_css, presentation_template
from .model import (
    ContentBlock,
    Presentation,
    PresentationModelError,
    Scene,
    Variant,
    stable_digest,
)

PUBLIC_METADATA_SCHEMA_VERSION = "schauwerk-stage-artifact-metadata.v1"
PRESENTER_SCHEMA_VERSION = "schauwerk-stage-presenter.v1"
_FIXED_ZIP_TIME = (1980, 1, 1, 0, 0, 0)


def variant_scenes(presentation: Presentation, variant: Variant) -> tuple[Scene, ...]:
    scene_map = presentation.scene_by_id
    return tuple(scene_map[item] for item in variant.scene_ids)


def _public_source_payload(
    presentation: Presentation, scenes: tuple[Scene, ...]
) -> list[dict[str, str]]:
    source_map = presentation.source_by_id
    source_ids = sorted(
        {source_id for scene in scenes for block in scene.blocks for source_id in block.source_ids}
    )
    return [
        {
            "id": source_id,
            "label": source_map[source_id].label,
            "revision": source_map[source_id].revision,
            "sha256": source_map[source_id].sha256,
        }
        for source_id in source_ids
    ]


def _visible_payload(scenes: tuple[Scene, ...]) -> list[dict[str, object]]:
    return [
        {
            "id": scene.identifier,
            "title": scene.title,
            "blocks": [
                {
                    "kind": block.kind,
                    "text": block.text,
                    "items": list(block.items),
                    "source_ids": list(block.source_ids),
                }
                for block in scene.blocks
            ],
        }
        for scene in scenes
    ]


def public_metadata(presentation: Presentation, variant: Variant) -> dict[str, object]:
    scenes = variant_scenes(presentation, variant)
    scene_order = [scene.identifier for scene in scenes]
    visible_payload = _visible_payload(scenes)
    public_sources = _public_source_payload(presentation, scenes)
    public_projection = {
        "presentation_id": presentation.presentation_id,
        "presentation_version": presentation.version,
        "variant_id": variant.identifier,
        "variant_title": variant.title,
        "audience": variant.audience,
        "source_revision": presentation.source_revision,
        "scene_order": scene_order,
        "visible": visible_payload,
        "sources": public_sources,
    }
    return {
        "schema_version": PUBLIC_METADATA_SCHEMA_VERSION,
        "presentation_id": presentation.presentation_id,
        "presentation_version": presentation.version,
        "variant_id": variant.identifier,
        "variant_title": variant.title,
        "audience": variant.audience,
        "source_revision": presentation.source_revision,
        "visual_grammar": GRAMMAR_SCHEMA_VERSION,
        "template": presentation_template().name,
        "scene_order": scene_order,
        "scene_order_sha256": stable_digest(scene_order),
        "visible_content_sha256": stable_digest(visible_payload),
        "public_sources": public_sources,
        "public_projection_sha256": stable_digest(public_projection),
    }


def _scene_source_labels(
    presentation: Presentation, scene: Scene, *, public_only: bool
) -> tuple[str, ...]:
    source_map = presentation.source_by_id
    identifiers: set[str] = set()
    if public_only:
        for block in scene.blocks:
            identifiers.update(block.source_ids)
    else:
        identifiers.update(scene.source_ids)
    return tuple(
        f"{source_map[item].label} ({source_map[item].revision})" for item in sorted(identifiers)
    )


def _block_html(block: ContentBlock) -> str:
    if block.kind == "bullets":
        return "<ul>" + "".join(f"<li>{html.escape(item)}</li>" for item in block.items) + "</ul>"
    if block.kind == "code":
        return f"<pre><code>{html.escape(block.text or '')}</code></pre>"
    css_class = "callout" if block.kind == "callout" else "paragraph"
    return f'<p class="{css_class}">{html.escape(block.text or "")}</p>'


def _metadata_tags(metadata: dict[str, object]) -> str:
    order = ",".join(metadata["scene_order"])
    return (
        f'<meta name="schauwerk-public-projection-sha256" '
        f'content="{metadata["public_projection_sha256"]}">'
        f'<meta name="schauwerk-visible-content-sha256" '
        f'content="{metadata["visible_content_sha256"]}">'
        f'<meta name="schauwerk-scene-order" content="{html.escape(order)}">'
        f'<meta name="schauwerk-source-revision" '
        f'content="{html.escape(str(metadata["source_revision"]))}">'
    )


def render_public_html(
    presentation: Presentation, variant: Variant, metadata: dict[str, object]
) -> str:
    scenes = variant_scenes(presentation, variant)
    style = (
        html_theme_css("presentation")
        + """
html {
  scroll-snap-type: y mandatory;
  scroll-behavior: smooth;
}
body { max-width: none; margin: 0; padding: 0; }
.slide {
  min-height: 100vh;
  box-sizing: border-box;
  padding: 7vh 8vw 6vh;
  scroll-snap-align: start;
  display: flex;
  flex-direction: column;
  justify-content: center;
  border: 0;
  border-bottom: 1px solid var(--sw-border);
  background: linear-gradient(150deg, #fff 0%, var(--sw-tint) 100%);
}
.slide h2 {
  font-size: clamp(2rem, 5vw, 4.8rem);
  line-height: 1.05;
  margin: 0 0 1.4rem;
  color: var(--sw-accent);
}
.slide p, .slide li {
  font-size: clamp(1.15rem, 2.2vw, 2rem);
  max-width: 46em;
}
.slide pre {
  font-size: clamp(.9rem, 1.6vw, 1.35rem);
  white-space: pre-wrap;
  background: #fff;
  border: 1px solid var(--sw-border);
  padding: 1rem;
}
.callout {
  border-left: .4rem solid var(--sw-accent);
  padding-left: 1rem;
  font-weight: 700;
}
.source { margin-top: auto; font-size: .85rem !important; color: #374151; }
.nav {
  position: fixed;
  right: 1rem;
  top: 1rem;
  z-index: 2;
  background: #fff;
  border: 1px solid var(--sw-border);
  padding: .45rem .7rem;
}
.nav a { margin-inline: .25rem; }
.meta { font-size: .8rem; color: #4B5563; }
@media print {
  .nav { display: none; }
  .slide { min-height: auto; height: 190mm; page-break-after: always; }
}
"""
    )
    navigation = "".join(
        f'<a href="#{html.escape(scene.identifier)}" aria-label="Folie {index}">{index}</a>'
        for index, scene in enumerate(scenes, start=1)
    )
    rendered_scenes = []
    for index, scene in enumerate(scenes, start=1):
        body = "".join(_block_html(block) for block in scene.blocks)
        sources = " · ".join(_scene_source_labels(presentation, scene, public_only=True))
        rendered_scenes.append(
            f'<section class="slide" id="{html.escape(scene.identifier)}" '
            f'data-scene-id="{html.escape(scene.identifier)}" aria-labelledby="title-{index}">'
            f'<p class="meta">{index}/{len(scenes)} · {html.escape(variant.audience)}</p>'
            f'<h2 id="title-{index}">{html.escape(scene.title)}</h2>{body}'
            f'<p class="source">Quellen: {html.escape(sources)} · Quellrevision: '
            f"{html.escape(presentation.source_revision)}</p></section>"
        )
    return (
        '<!doctype html>\n<html lang="de"><head><meta charset="utf-8">'
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        "<meta http-equiv='Content-Security-Policy' "
        "content=\"default-src 'none'; style-src 'unsafe-inline'; "
        "base-uri 'none'; form-action 'none'\">"
        f"{_metadata_tags(metadata)}<title>{html.escape(variant.title)}</title>"
        f"<style>{style}</style></head><body>"
        f'<nav class="nav" aria-label="Foliennavigation">{navigation}</nav>'
        f'<main data-presentation-id="{html.escape(presentation.presentation_id)}" '
        f'data-variant-id="{html.escape(variant.identifier)}">'
        f"{''.join(rendered_scenes)}</main></body></html>\n"
    )


def render_handout_html(
    presentation: Presentation, variant: Variant, metadata: dict[str, object]
) -> str:
    scenes = variant_scenes(presentation, variant)
    style = (
        html_theme_css("presentation")
        + """
body { max-width: 58rem; }
.handout {
  break-inside: avoid;
  border: 1px solid var(--sw-border);
  border-left: .4rem solid var(--sw-accent);
  padding: 1rem 1.4rem;
  margin: 1.2rem 0;
}
.handout h2 { margin-top: 0; }
.source { font-size: .82rem; color: #4B5563; }
.meta { font-size: .9rem; color: #4B5563; }
pre { white-space: pre-wrap; background: var(--sw-tint); padding: .8rem; }
@media print { .handout { page-break-inside: avoid; } }
"""
    )
    articles = []
    for index, scene in enumerate(scenes, start=1):
        sources = " · ".join(_scene_source_labels(presentation, scene, public_only=True))
        articles.append(
            f'<article class="handout" data-scene-id="{html.escape(scene.identifier)}">'
            f'<p class="meta">Folie {index}</p><h2>{html.escape(scene.title)}</h2>'
            f"{''.join(_block_html(block) for block in scene.blocks)}"
            f'<p class="source">Quellen: {html.escape(sources)}</p></article>'
        )
    return (
        '<!doctype html>\n<html lang="de"><head><meta charset="utf-8">'
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        "<meta http-equiv='Content-Security-Policy' "
        "content=\"default-src 'none'; style-src 'unsafe-inline'; "
        "base-uri 'none'; form-action 'none'\">"
        f"{_metadata_tags(metadata)}<title>{html.escape(variant.title)} · Handout</title>"
        f"<style>{style}</style></head><body><header><h1>{html.escape(variant.title)}</h1>"
        f'<p class="meta">Zielgruppe: {html.escape(variant.audience)} · Quellrevision: '
        f"{html.escape(presentation.source_revision)}</p></header><main>{''.join(articles)}</main>"
        '<footer><p class="meta">Öffentliches Handout ohne Sprecherhinweise '
        "und interne Zeitplanung.</p></footer></body></html>\n"
    )


def _wrap_text(text: str, *, font: str, size: float, width: float) -> list[str]:
    lines: list[str] = []
    for paragraph in text.splitlines() or [""]:
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        if stringWidth(current, font, size) > width:
            raise PresentationModelError(
                "text contains a token wider than the deterministic layout"
            )
        for word in words[1:]:
            if stringWidth(word, font, size) > width:
                raise PresentationModelError(
                    "text contains a token wider than the deterministic layout"
                )
            candidate = f"{current} {word}"
            if stringWidth(candidate, font, size) <= width:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def _pdf_block_height(block: ContentBlock, *, width: float) -> float:
    if block.kind == "bullets":
        return sum(
            max(1, len(_wrap_text(item, font="Helvetica", size=20, width=width - 26))) * 26 + 12
            for item in block.items
        )
    font = (
        "Courier"
        if block.kind == "code"
        else "Helvetica-Bold"
        if block.kind == "callout"
        else "Helvetica"
    )
    size = 15 if block.kind == "code" else 21
    indent = 18 if block.kind == "callout" else 0
    lines = _wrap_text(block.text or "", font=font, size=size, width=width - indent)
    return len(lines) * size * 1.35 + 16


def _draw_pdf_block(
    canvas: Canvas, block: ContentBlock, *, x: float, y: float, width: float
) -> float:
    if block.kind == "bullets":
        for item in block.items:
            lines = _wrap_text(item, font="Helvetica", size=20, width=width - 26)
            canvas.setFont("Helvetica", 20)
            canvas.drawString(x, y, "-")
            for line_index, line in enumerate(lines):
                canvas.drawString(x + 24, y - line_index * 26, line)
            y -= max(1, len(lines)) * 26 + 12
        return y
    font = (
        "Courier"
        if block.kind == "code"
        else "Helvetica-Bold"
        if block.kind == "callout"
        else "Helvetica"
    )
    size = 15 if block.kind == "code" else 21
    indent = 18 if block.kind == "callout" else 0
    if block.kind == "callout":
        canvas.setLineWidth(4)
        canvas.line(x, y + 6, x, y - 70)
    canvas.setFont(font, size)
    lines = _wrap_text(block.text or "", font=font, size=size, width=width - indent)
    for line in lines:
        canvas.drawString(x + indent, y, line)
        y -= size * 1.35
    return y - 16


def _validate_pdf_text(presentation: Presentation, variant: Variant) -> None:
    scenes = variant_scenes(presentation, variant)
    values = [variant.title, presentation.source_revision]
    for scene in scenes:
        values.extend(scene.visible_strings())
        values.extend(_scene_source_labels(presentation, scene, public_only=True))
    for value in values:
        try:
            value.encode("cp1252")
        except UnicodeEncodeError as exc:
            raise PresentationModelError(
                "public text contains a character unsupported by the deterministic PDF font"
            ) from exc


def render_pdf(
    path: Path, presentation: Presentation, variant: Variant, metadata: dict[str, object]
) -> None:
    _validate_pdf_text(presentation, variant)
    pagesize = landscape(A4)
    width, height = pagesize
    canvas = Canvas(str(path), pagesize=pagesize, invariant=1, pageCompression=0)
    canvas.setAuthor("Schauwerk")
    canvas.setCreator("Schauwerk SW-012")
    canvas.setTitle(variant.title)
    canvas.setSubject(
        f"public_projection_sha256={metadata['public_projection_sha256']};"
        f"visible_content_sha256={metadata['visible_content_sha256']};"
        f"scene_order_sha256={metadata['scene_order_sha256']}"
    )
    scenes = variant_scenes(presentation, variant)
    for index, scene in enumerate(scenes, start=1):
        canvas.setFillColorRGB(0.08, 0.13, 0.24)
        canvas.rect(0, height - 52, width, 52, fill=1, stroke=0)
        canvas.setFillColorRGB(1, 1, 1)
        canvas.setFont("Helvetica", 11)
        canvas.drawString(34, height - 33, variant.title)
        canvas.drawRightString(width - 34, height - 33, f"{index}/{len(scenes)}")
        canvas.setFillColorRGB(0.08, 0.13, 0.24)
        canvas.setFont("Helvetica-Bold", 30)
        title_lines = _wrap_text(scene.title, font="Helvetica-Bold", size=30, width=width - 80)
        y = height - 92
        for line in title_lines:
            canvas.drawString(40, y, line)
            y -= 38
        y -= 8
        content_width = width - 96
        required_height = sum(
            _pdf_block_height(block, width=content_width) for block in scene.blocks
        )
        if y - required_height < 85:
            raise PresentationModelError(
                f"scene {scene.identifier} exceeds the deterministic PDF layout"
            )
        for block in scene.blocks:
            y = _draw_pdf_block(canvas, block, x=48, y=y, width=content_width)
        sources = " · ".join(_scene_source_labels(presentation, scene, public_only=True))
        source_footer = f"Quellen: {sources}"
        revision_footer = f"Revision {presentation.source_revision}"
        if stringWidth(source_footer, "Helvetica", 8) > width - 300:
            raise PresentationModelError(
                f"scene {scene.identifier} source footer exceeds the deterministic PDF layout"
            )
        if stringWidth(revision_footer, "Helvetica", 8) > 220:
            raise PresentationModelError("source revision exceeds the deterministic PDF layout")
        canvas.setFont("Helvetica", 8)
        canvas.setFillColorRGB(0.25, 0.28, 0.32)
        canvas.drawString(40, 32, source_footer)
        canvas.drawRightString(width - 40, 32, revision_footer)
        canvas.showPage()
    canvas.save()


def _add_pptx_text(
    slide,
    x,
    y,
    width,
    height,
    text,
    *,
    size,
    bold=False,
    font="Arial",
    color=(20, 33, 61),
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return shape


def _pptx_line_count(text: str, *, size: float, width_inches: float) -> int:
    return len(_wrap_text(text, font="Helvetica", size=size, width=width_inches * 72))


def _pptx_block_height(block: ContentBlock) -> float:
    if block.kind == "bullets":
        line_count = sum(_pptx_line_count(item, size=22, width_inches=10.7) for item in block.items)
        return 0.35 + line_count * 0.5
    text = block.text or ""
    size = 16 if block.kind == "code" else 23
    line_count = _pptx_line_count(text, size=size, width_inches=11.2)
    line_height = 0.34 if block.kind == "code" else 0.46
    return 0.35 + line_count * line_height


def _render_pptx_block(slide, block: ContentBlock, *, y: float) -> float:
    height = _pptx_block_height(block)
    if y + height > 6.55:
        raise PresentationModelError("scene exceeds the deterministic PowerPoint layout")
    if block.kind == "bullets":
        text = "\n".join(f"• {item}" for item in block.items)
        _add_pptx_text(slide, 1.0, y, 11.2, height, text, size=22)
        return y + height + 0.18
    text = block.text or ""
    size = 16 if block.kind == "code" else 23
    font = "Courier New" if block.kind == "code" else "Arial"
    bold = block.kind == "callout"
    shape = _add_pptx_text(slide, 1.0, y, 11.2, height, text, size=size, bold=bold, font=font)
    if block.kind == "callout":
        shape.line.color.rgb = RGBColor(20, 33, 61)
        shape.line.width = Pt(2)
    return y + height + 0.18


def _normalize_pptx_zip(source: Path, destination: Path) -> None:
    with ZipFile(source, "r") as archive:
        members = [(name, archive.read(name)) for name in sorted(archive.namelist())]
    with ZipFile(destination, "w", compression=ZIP_DEFLATED, compresslevel=9) as archive:
        for name, payload in members:
            info = ZipInfo(filename=name, date_time=_FIXED_ZIP_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o100644 << 16
            archive.writestr(info, payload)


def render_pptx(
    path: Path, presentation: Presentation, variant: Variant, metadata: dict[str, object]
) -> None:
    deck = PptxPresentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    properties = deck.core_properties
    properties.title = variant.title
    properties.subject = (
        f"public_projection_sha256={metadata['public_projection_sha256']};"
        f"visible_content_sha256={metadata['visible_content_sha256']}"
    )
    properties.author = "Schauwerk"
    properties.language = "de-DE"
    properties.keywords = ",".join(metadata["scene_order"])
    properties.comments = f"source_revision={presentation.source_revision}"
    properties.created = datetime(2000, 1, 1)
    properties.modified = datetime(2000, 1, 1)
    blank_layout = deck.slide_layouts[6]
    scenes = variant_scenes(presentation, variant)
    if _pptx_line_count(variant.title, size=11, width_inches=10.8) > 1:
        raise PresentationModelError(
            "variant title exceeds the deterministic PowerPoint header layout"
        )
    for index, scene in enumerate(scenes, start=1):
        if _pptx_line_count(scene.title, size=31, width_inches=11.9) > 2:
            raise PresentationModelError(
                f"scene {scene.identifier} title exceeds the deterministic PowerPoint layout"
            )
        slide = deck.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(237, 244, 255)
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333333), Inches(0.55)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(20, 33, 61)
        banner.line.fill.background()
        _add_pptx_text(
            slide,
            0.55,
            0.1,
            10.8,
            0.35,
            variant.title,
            size=11,
            bold=True,
            color=(255, 255, 255),
        )
        count = _add_pptx_text(
            slide,
            11.65,
            0.1,
            1.0,
            0.35,
            f"{index}/{len(scenes)}",
            size=11,
            bold=True,
            color=(255, 255, 255),
        )
        count.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _add_pptx_text(slide, 0.75, 0.82, 11.9, 1.05, scene.title, size=31, bold=True)
        y = 1.9
        for block in scene.blocks:
            try:
                y = _render_pptx_block(slide, block, y=y)
            except PresentationModelError as exc:
                raise PresentationModelError(
                    f"scene {scene.identifier} exceeds the deterministic PowerPoint layout"
                ) from exc
        sources = " · ".join(_scene_source_labels(presentation, scene, public_only=True))
        footer_text = f"Quellen: {sources} · Revision {presentation.source_revision}"
        if _pptx_line_count(footer_text, size=8, width_inches=11.9) > 1:
            raise PresentationModelError(
                f"scene {scene.identifier} footer exceeds the deterministic PowerPoint layout"
            )
        footer = _add_pptx_text(
            slide,
            0.75,
            6.92,
            11.9,
            0.32,
            footer_text,
            size=8,
        )
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    fd, temporary_name = tempfile.mkstemp(prefix="schauwerk-pptx-", suffix=".pptx", dir=path.parent)
    os.close(fd)
    temporary = Path(temporary_name)
    try:
        deck.save(temporary)
        _normalize_pptx_zip(temporary, path)
    finally:
        temporary.unlink(missing_ok=True)


def presenter_payload(presentation: Presentation, variant: Variant) -> dict[str, object]:
    scenes = variant_scenes(presentation, variant)
    source_map = presentation.source_by_id
    return {
        "schema_version": PRESENTER_SCHEMA_VERSION,
        "presentation_id": presentation.presentation_id,
        "presentation_version": presentation.version,
        "variant_id": variant.identifier,
        "variant_title": variant.title,
        "audience": variant.audience,
        "source_revision": presentation.source_revision,
        "model_digest": presentation.model_digest,
        "planned_duration_seconds": variant.planned_duration_seconds,
        "scene_order": [scene.identifier for scene in scenes],
        "scenes": [
            {
                "id": scene.identifier,
                "title": scene.title,
                "duration_seconds": scene.duration_seconds,
                "speaker_notes": list(scene.speaker_notes),
                "sources": [
                    {
                        "id": source_id,
                        "label": source_map[source_id].label,
                        "revision": source_map[source_id].revision,
                        "visibility": source_map[source_id].visibility,
                        "sha256": source_map[source_id].sha256,
                    }
                    for source_id in scene.source_ids
                ],
            }
            for scene in scenes
        ],
        "boundaries": {
            "public_distribution_allowed": False,
            "contains_speaker_notes": True,
            "contains_timing": True,
            "contains_artifact_paths": False,
            "network_dependencies": False,
        },
    }


def render_presenter_html(payload: dict[str, object]) -> str:
    scene_cards = []
    for index, scene in enumerate(payload["scenes"], start=1):
        notes = "".join(f"<li>{html.escape(note)}</li>" for note in scene["speaker_notes"])
        sources = " · ".join(
            f"{item['label']} ({item['revision']}, {item['visibility']})"
            for item in scene["sources"]
        )
        scene_cards.append(
            f'<section class="scene" data-scene-id="{html.escape(scene["id"])}">'
            f'<p class="meta">{index}/{len(payload["scenes"])} · '
            f"{scene['duration_seconds']} Sekunden</p><h2>{html.escape(scene['title'])}</h2>"
            f"<h3>Sprecherhinweise</h3><ul>{notes}</ul>"
            f'<p class="sources">Quellen: {html.escape(sources)}</p></section>'
        )
    style = """
:root {
  font-family: system-ui, sans-serif;
  color: #181818;
  background: #f5f7fa;
}
body { max-width: 72rem; margin: 0 auto; padding: 2rem; }
.warning {
  border: 3px solid #8a2c0d;
  background: #fff1eb;
  padding: 1rem;
  font-weight: 700;
}
.scene {
  background: #fff;
  border: 1px solid #64748b;
  border-left: .45rem solid #14213d;
  padding: 1rem 1.4rem;
  margin: 1.2rem 0;
}
.meta, .sources { color: #475569; font-size: .9rem; }
h1, h2 { color: #14213d; }
"""
    return (
        '<!doctype html>\n<html lang="de"><head><meta charset="utf-8">'
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        "<meta http-equiv='Content-Security-Policy' "
        "content=\"default-src 'none'; style-src 'unsafe-inline'; "
        "base-uri 'none'; form-action 'none'\">"
        f'<meta name="schauwerk-model-digest" content="{payload["model_digest"]}">'
        f"<title>{html.escape(payload['variant_title'])} · Presenter</title>"
        f"<style>{style}</style></head><body><header><h1>{html.escape(payload['variant_title'])}</h1>"
        '<p class="warning">Internes Presenter-Paket: nicht projizieren oder veröffentlichen.</p>'
        f"<p>Gesamtzeit: {payload['planned_duration_seconds']} Sekunden · Zielgruppe: "
        f"{html.escape(payload['audience'])}</p></header><main>{''.join(scene_cards)}</main>"
        "</body></html>\n"
    )


def sha256_file(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()
